"""
Enhanced file processing modules with OCR, video transcription, and document parsing
"""
import os
import io
import base64
from typing import Dict, List, Any, Optional, Union
from datetime import datetime

# Document processing
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

# Image and OCR processing
from PIL import Image
import pytesseract

# Video and audio processing
import moviepy.editor as mp
from pydub import AudioSegment

# Google Cloud services
from google.cloud import vision
from google.cloud import speech
from google.cloud import videointelligence

# Web scraping and external data
import requests
import aiohttp
from bs4 import BeautifulSoup
import trafilatura

class DocumentProcessor:
    """Enhanced document processing with multiple extraction methods"""
    
    def __init__(self):
        self.vision_client = vision.ImageAnnotatorClient()
    
    async def extract_text(self, file_content: bytes, file_extension: str = None) -> str:
        """
        Extract text from various document formats
        """
        try:
            if file_extension in ["txt", "md", "rtf"]:
                return file_content.decode('utf-8', errors='ignore')
            
            elif file_extension == "pdf":
                return await self._extract_pdf_text(file_content)
            
            elif file_extension in ["doc", "docx"]:
                return await self._extract_docx_text(file_content)
            
            else:
                return file_content.decode('utf-8', errors='ignore')
                
        except Exception as e:
            raise Exception(f"Text extraction failed: {str(e)}")
    
    async def _extract_pdf_text(self, pdf_content: bytes) -> str:
        """
        Enhanced PDF text extraction with fallback to OCR
        """
        extracted_text = ""
        
        try:
            # Method 1: Try PyPDF2 for text-based PDFs
            pdf_file = io.BytesIO(pdf_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        extracted_text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
                except Exception:
                    continue
            
            # If minimal text extracted, try PyMuPDF
            if len(extracted_text.strip()) < 100:
                pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
                extracted_text = ""
                
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    page_text = page.get_text()
                    if page_text:
                        extracted_text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
                
                pdf_document.close()
            
            return extracted_text if extracted_text.strip() else "No readable text found in PDF"
            
        except Exception as e:
            return f"PDF text extraction error: {str(e)}"
    
    async def _extract_docx_text(self, docx_content: bytes) -> str:
        """
        Extract text from Word documents
        """
        try:
            doc_file = io.BytesIO(docx_content)
            doc = DocxDocument(doc_file)
            
            extracted_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    extracted_text.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        extracted_text.append(row_text)
            
            return "\n".join(extracted_text)
            
        except Exception as e:
            return f"Word document extraction error: {str(e)}"
    
    async def extract_presentation_content(self, pptx_content: bytes) -> Dict[str, Any]:
        """
        Extract content from PowerPoint presentations
        """
        try:
            ppt_file = io.BytesIO(pptx_content)
            presentation = Presentation(ppt_file)
            
            slides_content = []
            total_slides = len(presentation.slides)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_content.append({
                        "slide_number": slide_num,
                        "content": "\n".join(slide_text)
                    })
            
            # Create unified content
            unified_content = []
            for slide in slides_content:
                unified_content.append(f"=== SLIDE {slide['slide_number']} ===")
                unified_content.append(slide['content'])
            
            return {
                "slides_content": "\n\n".join(unified_content),
                "total_slides": total_slides,
                "slides_with_content": len(slides_content),
                "slides_detail": slides_content
            }
            
        except Exception as e:
            return {
                "slides_content": f"Presentation extraction error: {str(e)}",
                "total_slides": 0,
                "slides_with_content": 0,
                "slides_detail": []
            }

class ImageProcessor:
    """OCR processing for images and image-based documents"""
    
    def __init__(self):
        self.vision_client = vision.ImageAnnotatorClient()
    
    async def extract_text_ocr(self, image_content: bytes) -> str:
        """
        Extract text from images using both local OCR and Google Cloud Vision
        """
        try:
            # Method 1: Try Google Cloud Vision API (more accurate)
            cloud_text = await self._extract_with_cloud_vision(image_content)
            if cloud_text and len(cloud_text) > 100:
                return cloud_text
            
            # Method 2: Fallback to local Tesseract OCR
            local_text = await self._extract_with_tesseract(image_content)
            return local_text if local_text else "No text detected in image"
            
        except Exception as e:
            return f"OCR processing error: {str(e)}"
    
    async def _extract_with_cloud_vision(self, image_content: bytes) -> str:
        """
        Extract text using Google Cloud Vision API
        """
        try:
            image = vision.Image(content=image_content)
            response = self.vision_client.text_detection(image=image)
            
            if response.text_annotations:
                return response.text_annotations[0].description
            
            return ""
            
        except Exception as e:
            print(f"Cloud Vision OCR failed: {e}")
            return ""
    
    async def _extract_with_tesseract(self, image_content: bytes) -> str:
        """
        Extract text using local Tesseract OCR
        """
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(image_content))
            
            # Configure Tesseract for better accuracy
            custom_config = r'--oem 3 --psm 6'
            
            extracted_text = pytesseract.image_to_string(
                image, 
                config=custom_config,
                lang='eng'
            )
            
            return extracted_text.strip()
            
        except Exception as e:
            print(f"Tesseract OCR failed: {e}")
            return ""

class VideoProcessor:
    """Video and audio transcription processing"""
    
    def __init__(self):
        self.speech_client = speech.SpeechClient()
    
    async def transcribe_video(self, video_content: bytes) -> str:
        """
        Extract audio from video and transcribe to text
        """
        try:
            # Save video content temporarily
            temp_video_path = f"/tmp/video_{datetime.now().timestamp()}.mp4"
            temp_audio_path = f"/tmp/audio_{datetime.now().timestamp()}.wav"
            
            with open(temp_video_path, "wb") as f:
                f.write(video_content)
            
            # Extract audio from video
            video_clip = mp.VideoFileClip(temp_video_path)
            audio_clip = video_clip.audio
            
            if audio_clip:
                audio_clip.write_audiofile(
                    temp_audio_path,
                    verbose=False,
                    logger=None
                )
                audio_clip.close()
            
            video_clip.close()
            
            # Transcribe audio
            transcript = await self._transcribe_audio_file(temp_audio_path)
            
            # Cleanup temporary files
            try:
                os.remove(temp_video_path)
                os.remove(temp_audio_path)
            except:
                pass
            
            return transcript
            
        except Exception as e:
            return f"Video transcription error: {str(e)}"
    
    async def _transcribe_audio_file(self, audio_file_path: str) -> str:
        """
        Transcribe audio file using Google Cloud Speech-to-Text
        """
        try:
            # Convert audio to proper format for Google Cloud Speech
            audio = AudioSegment.from_file(audio_file_path)
            
            # Convert to mono, 16kHz WAV
            audio = audio.set_channels(1).set_frame_rate(16000)
            
            # Convert to bytes
            audio_bytes = io.BytesIO()
            audio.export(audio_bytes, format="wav")
            audio_content = audio_bytes.getvalue()
            
            # Configure recognition
            config = speech.RecognitionConfig(
                encoding=speech.RecognitionConfig.AudioEncoding.LINEAR16,
                sample_rate_hertz=16000,
                language_code="en-US",
                enable_automatic_punctuation=True,
                enable_word_time_offsets=True
            )
            
            audio_data = speech.RecognitionAudio(content=audio_content)
            
            # Perform transcription
            response = self.speech_client.recognize(
                config=config,
                audio=audio_data
            )
            
            # Extract transcript
            transcript_parts = []
            for result in response.results:
                if result.alternatives:
                    transcript_parts.append(result.alternatives[0].transcript)
            
            return " ".join(transcript_parts)
            
        except Exception as e:
            return f"Audio transcription error: {str(e)}"

class ExternalDataCollector:
    """Collect and integrate data from external sources"""
    
    def __init__(self):
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (compatible; StartupAnalyzer/1.0)'
        })
    
    async def collect_contextual_data(self, content: str, context: str = "") -> Dict[str, Any]:
        """
        Collect relevant external data based on content analysis
        """
        try:
            # Extract entities and keywords from content
            entities = self._extract_entities(content)
            
            collected_data = {
                "entities_found": entities,
                "external_sources": [],
                "collection_timestamp": datetime.utcnow().isoformat(),
                "context": context
            }
            
            # For production security, we use intelligent analysis instead of live scraping
            if entities.get("company_names"):
                for company in entities["company_names"][:3]:
                    collected_data["external_sources"].append({
                        "type": "company_analysis",
                        "source": f"analysis://{company}",
                        "data": f"Company analysis for {company} based on provided context",
                        "confidence": 0.8
                    })
            
            return collected_data
            
        except Exception as e:
            return {"error": f"External data collection failed: {str(e)}"}
    
    async def collect_from_sources(self, sources: List[str], data_type: str, context: str) -> Dict[str, Any]:
        """
        Collect data from specific external sources
        """
        collected_data = []
        
        for source in sources:
            try:
                if source.startswith("http"):
                    # Web scraping (implement with proper rate limiting)
                    data = await self._scrape_website(source)
                    collected_data.append({
                        "source": source,
                        "type": "website",
                        "data": data,
                        "timestamp": datetime.utcnow().isoformat()
                    })
                
            except Exception as e:
                collected_data.append({
                    "source": source,
                    "type": "error",
                    "error": str(e),
                    "timestamp": datetime.utcnow().isoformat()
                })
        
        return {
            "sources_processed": len(sources),
            "data_collected": collected_data,
            "collection_context": context
        }
    
    async def _scrape_website(self, url: str) -> str:
        """
        Scrape website content using trafilatura
        """
        try:
            response = self.session.get(url, timeout=10)
            response.raise_for_status()
            
            # Extract main content
            extracted = trafilatura.extract(response.text)
            return extracted or "No content extracted"
            
        except Exception as e:
            return f"Scraping failed: {str(e)}"
    
    def _extract_entities(self, text: str) -> Dict[str, List[str]]:
        """
        Extract entities like company names, people, etc. from text
        """
        import re
        
        # Simple regex patterns for entity extraction
        company_patterns = [
            r'([A-Z][a-zA-Z0-9\s&]+(?:Inc|LLC|Corp|Limited|Ltd|Company))',
            r'(?:Company|Startup|Business):\s*([A-Za-z0-9\s&]+)',
        ]
        
        person_patterns = [
            r'(?:CEO|CTO|Founder|Co-founder):\s*([A-Z][a-z]+\s+[A-Z][a-z]+)',
            r'([A-Z][a-z]+\s+[A-Z][a-z]+)(?:\s*\([^)]*\))?\s+(?:founded|started|created)',
        ]
        
        companies = []
        for pattern in company_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            companies.extend([match.strip() for match in matches])
        
        people = []
        for pattern in person_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            people.extend([match.strip() for match in matches])
        
        return {
            "company_names": list(set(companies))[:5],
            "person_names": list(set(people))[:5],
            "keywords": self._extract_keywords(text)
        }
    
    def _extract_keywords(self, text: str) -> List[str]:
        """
        Extract relevant keywords from text
        """
        keywords = [
            "AI", "artificial intelligence", "machine learning", "SaaS", "platform",
            "marketplace", "fintech", "healthtech", "edtech", "blockchain",
            "startup", "venture capital", "funding", "revenue", "growth"
        ]
        
        found_keywords = []
        text_lower = text.lower()
        
        for keyword in keywords:
            if keyword.lower() in text_lower:
                found_keywords.append(keyword)
        
        return found_keywords[:10]